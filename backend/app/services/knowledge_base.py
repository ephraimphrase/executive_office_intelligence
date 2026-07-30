import io
import logging
import math
import uuid as uuid_lib
from datetime import datetime, timezone

from app.integrations.openai_client import get_openai_client

logger = logging.getLogger(__name__)


def _cosine_similarity(a: list, b: list) -> float:
    if not a or not b or len(a) != len(b):
        return -1.0
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = math.sqrt(sum(x * x for x in a))
    norm_b = math.sqrt(sum(y * y for y in b))
    if norm_a == 0 or norm_b == 0:
        return -1.0
    return dot / (norm_a * norm_b)


class KnowledgeBaseService:
    """Semantic search backed by pgvector — document embeddings live directly
    on the Document row (app/models/document.py: Document.embedding) rather
    than in a separate external search index."""

    def __init__(self):
        self.openai_client = get_openai_client()

    async def index_document(self, doc_id: str, name: str, content_text: str, file_type: str, metadata: dict, db) -> bool:
        """Embed the document's content and store the vector on its Document row."""
        if db is None:
            return False

        vector = await self.openai_client.embed(content_text[:2000])

        from app.models.document import Document

        try:
            doc_uuid = uuid_lib.UUID(str(doc_id))
        except ValueError:
            return False

        doc = await db.get(Document, doc_uuid)
        if not doc:
            return False
        doc.embedding = vector
        await db.commit()
        return True

    async def search(self, query: str, top_k: int, filters: dict, db) -> list:
        """Semantic search over indexed documents. Uses real pgvector cosine-
        distance ordering on Postgres; falls back to a Python-side cosine scan
        on SQLite (dev), where pgvector's column type isn't available."""
        if db is None:
            return []

        from sqlalchemy import select

        from app.database import _is_sqlite
        from app.models.document import Document

        query_vector = await self.openai_client.embed(query)

        if _is_sqlite:
            result = await db.execute(select(Document).where(Document.embedding.isnot(None)))
            scored = [
                (_cosine_similarity(query_vector, doc.embedding), doc)
                for doc in result.scalars().all()
            ]
            scored.sort(key=lambda pair: pair[0], reverse=True)
            return [doc for score, doc in scored[:top_k] if score > -1.0]

        result = await db.execute(
            select(Document)
            .where(Document.embedding.isnot(None))
            .order_by(Document.embedding.cosine_distance(query_vector))
            .limit(top_k)
        )
        return result.scalars().all()

    async def extract_text_from_file(self, file_bytes: bytes, file_type: str) -> str:
        """Extract plain text from Word/PDF/Excel/PowerPoint/other files for indexing."""
        if not file_bytes:
            return ""
        file_type = (file_type or "OTHER").upper()
        try:
            if file_type == "WORD":
                return self._extract_docx(file_bytes)
            if file_type == "PDF":
                return self._extract_pdf(file_bytes)
            if file_type == "EXCEL":
                return self._extract_xlsx(file_bytes)
            if file_type == "POWERPOINT":
                return self._extract_pptx(file_bytes)
            return file_bytes.decode("utf-8", errors="ignore")
        except Exception as e:
            logger.warning(f"Text extraction failed for file_type={file_type}: {e}")
            return ""

    def _extract_docx(self, file_bytes: bytes) -> str:
        import docx
        document = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in document.paragraphs if p.text)

    def _extract_pdf(self, file_bytes: bytes) -> str:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    def _extract_xlsx(self, file_bytes: bytes) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        lines = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                values = [str(v) for v in row if v is not None]
                if values:
                    lines.append(" ".join(values))
        return "\n".join(lines)

    def _extract_pptx(self, file_bytes: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines)

    def _classify_document(self, name: str, text: str) -> dict:
        """Keyword-based classification so board papers/policies/contracts
        pulled in from OneDrive are auto-flagged instead of staying
        uncategorized until someone manually edits them."""
        haystack = f"{name or ''} {(text or '')[:1000]}".lower()

        is_board_paper = any(kw in haystack for kw in (
            "board paper", "board memo", "board resolution", "for board approval", "board pack",
        ))
        is_policy = any(kw in haystack for kw in (
            "policy", "standard operating procedure", " sop ", "guideline",
        ))
        is_contract = any(kw in haystack for kw in (
            "contract", "agreement", "memorandum of understanding", " mou ", "nda",
        ))

        if is_board_paper:
            category, subcategory = "BOARD", "BOARD_PAPER"
        elif is_contract:
            category, subcategory = "LEGAL", "CONTRACT"
        elif is_policy:
            category, subcategory = "POLICY", "POLICY"
        else:
            category, subcategory = "GENERAL", "UNCATEGORIZED"

        return {
            "is_board_paper": is_board_paper,
            "is_policy": is_policy,
            "is_contract": is_contract,
            "category": category,
            "subcategory": subcategory,
        }

    def _infer_file_type(self, name: str, mime: str = "") -> str:
        lower = (name or "").lower()
        mime = mime or ""
        if lower.endswith(".docx") or "wordprocessingml" in mime:
            return "WORD"
        if lower.endswith(".pdf") or "pdf" in mime:
            return "PDF"
        if lower.endswith(".xlsx") or "spreadsheetml.sheet" in mime:
            return "EXCEL"
        if lower.endswith(".pptx") or "presentationml" in mime:
            return "POWERPOINT"
        return "OTHER"

    async def sync_onedrive(self, db) -> int:
        """List files from OneDrive, download+extract new ones, and index them."""
        logger.info("Syncing knowledge base from onedrive...")
        if db is None:
            return 0

        from sqlalchemy import select

        from app.config import get_settings
        from app.integrations.microsoft_graph import MicrosoftGraphClient
        from app.models.document import Document

        settings = get_settings()
        graph = MicrosoftGraphClient()
        items = await graph.list_drive_items(settings.gvp_email, "/")

        count = 0
        for item in items:
            onedrive_id = item.get("id")
            if not onedrive_id:
                continue

            existing = await db.execute(select(Document).where(Document.onedrive_id == onedrive_id))
            if existing.scalars().first():
                continue

            name = item.get("name", "Untitled")
            mime = (item.get("file") or {}).get("mimeType", "")
            file_type = self._infer_file_type(name, mime)

            file_bytes = await graph.download_file(settings.gvp_email, onedrive_id)
            text = await self.extract_text_from_file(file_bytes, file_type)
            classification = self._classify_document(name, text)

            doc = Document(
                name=name,
                file_type=file_type,
                onedrive_id=onedrive_id,
                size_bytes=item.get("size"),
                ai_summary=text[:500] if text else None,
                indexed_at=datetime.now(timezone.utc),
                **classification,
            )
            db.add(doc)
            await db.flush()

            await self.index_document(str(doc.id), name, text, file_type, {}, db)
            count += 1

        await db.commit()
        return count

    async def get_document_summary(self, doc_id: str, db) -> str:
        from app.models.document import Document
        doc = await db.get(Document, doc_id)
        if not doc:
            return ""
        return doc.ai_summary or "No summary available."
