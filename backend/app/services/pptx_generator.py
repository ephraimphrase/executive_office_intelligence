import io
import logging

logger = logging.getLogger(__name__)


def _fmt(value) -> str:
    if not value:
        return ""
    if hasattr(value, "strftime"):
        return value.strftime("%Y-%m-%d %H:%M")
    return str(value)


class PowerPointGeneratorService:
    """Generates real .pptx board/briefing decks — the counterpart to
    WordGeneratorService for when the GVP needs a slide deck rather than a
    document (board packs, executive briefing decks)."""

    def _new_deck(self, title: str, subtitle: str = ""):
        from pptx import Presentation

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        return prs

    def _add_bullet_slide(self, prs, title: str, bullets: list[str]):
        bullet_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        if not bullets:
            body.text = "None"
            return slide

        body.text = bullets[0]
        for bullet in bullets[1:]:
            p = body.add_paragraph()
            p.text = bullet
        return slide

    def _add_table_slide(self, prs, title: str, headers: list[str], rows: list[list[str]]):
        blank_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = title

        from pptx.util import Inches

        n_rows = max(len(rows), 0) + 1
        n_cols = len(headers)
        table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * n_rows))
        table = table_shape.table
        for col, header in enumerate(headers):
            table.cell(0, col).text = header
        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                table.cell(r, c).text = str(value)
        return slide

    def _to_bytes(self, prs) -> bytes:
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    async def generate_briefing_deck(self, briefing_pack: dict, db) -> bytes:
        """Turns the Executive Briefing Pack (events, tasks, risks, decisions,
        priorities) into a board-ready slide deck."""
        date_str = str(briefing_pack.get("date", ""))
        prs = self._new_deck("Executive Briefing", date_str)

        priorities = briefing_pack.get("priorities") or []
        if priorities:
            self._add_bullet_slide(prs, "Today's Priorities", [str(p) for p in priorities])

        events = briefing_pack.get("events") or []
        if events:
            rows = [
                [_fmt(e.get("start_datetime")), e.get("title", ""), e.get("location") or ""]
                for e in events
            ]
            self._add_table_slide(prs, "Today's Schedule", ["Time", "Meeting", "Location"], rows)

        risks = briefing_pack.get("risks") or []
        if risks:
            self._add_bullet_slide(
                prs, "Open Risks",
                [f"[{r.get('severity', '')}] {r.get('description', '')}" for r in risks],
            )

        decisions = briefing_pack.get("decisions") or []
        if decisions:
            self._add_bullet_slide(
                prs, "Pending Decisions",
                [d.get("description", "") for d in decisions],
            )

        talking_points = briefing_pack.get("talking_points") or []
        if talking_points:
            self._add_bullet_slide(prs, "Talking Points", [str(t) for t in talking_points])

        return self._to_bytes(prs)

    async def generate_meeting_deck(self, meeting: dict, db) -> bytes:
        """Board pack for a single meeting: agenda + participants + prep notes."""
        prs = self._new_deck(meeting.get("title", "Meeting"), _fmt(meeting.get("meeting_date")))

        agenda = meeting.get("agenda") or []
        agenda_items = [a.get("item", str(a)) if isinstance(a, dict) else str(a) for a in agenda]
        self._add_bullet_slide(prs, "Agenda", agenda_items)

        participants = meeting.get("participants") or []
        if participants:
            self._add_bullet_slide(prs, "Participants", [str(p) for p in participants])

        talking_points = meeting.get("talking_points") or []
        if talking_points:
            self._add_bullet_slide(prs, "Talking Points", [str(t) for t in talking_points])

        return self._to_bytes(prs)
